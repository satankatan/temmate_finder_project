# -*- coding: utf-8 -*-
"""Генерация презентации по проекту Teammate Finder."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT = Path(__file__).resolve().parent.parent
OUTPUT = PROJECT / "presentation" / "teammate_finder_presentation.pptx"
GRAPHS = PROJECT / "results" / "custom_model"

# Палитра как в teammate_finder.html
BG = RGBColor(0x0B, 0x0F, 0x1A)
SURFACE = RGBColor(0x12, 0x18, 0x29)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT2 = RGBColor(0x00, 0xCE, 0xC9)
TEXT = RGBColor(0xE8, 0xED, 0xF7)
MUTED = RGBColor(0x8B, 0x97, 0xB3)
SUCCESS = RGBColor(0x00, 0xB8, 0x94)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE = (
    "Нейросетевая система семантического поиска\n"
    "идеального тиммейта"
)
SUBTITLE = "на основе собственной модели и размеченных игровых данных"


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text: str = "Teammate Finder · ЮФУ · 2025"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9), Inches(0.35))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def style_title(tf, size=36, bold=True, color=TEXT):
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def add_bullets(text_frame, items: list[str], size=18, color=TEXT, spacing=8):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)


def add_accent_bar(slide, left=0, top=0, width=0.12, height=7.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_card(slide, left, top, width, height, title: str, body: str, accent=ACCENT):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = SURFACE
    rect.line.color.rgb = RGBColor(0x2A, 0x35, 0x50)

    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(6)


def slide_title_only(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(8.8), Inches(1.2))
    tf = tb.text_frame
    tf.text = title
    style_title(tf, size=40)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = ACCENT2
        p2.space_before = Pt(12)
    add_footer(slide)
    return slide


def slide_content(prs: Presentation, title: str, bullets: list[str], subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.text = title
    style_title(tf, size=32)

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(8.8), Inches(0.5))
        st.text_frame.text = subtitle
        st.text_frame.paragraphs[0].font.size = Pt(14)
        st.text_frame.paragraphs[0].font.color.rgb = MUTED

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7 if subtitle else 1.4), Inches(8.8), Inches(5.2))
    add_bullets(body.text_frame, bullets, size=17)
    add_footer(slide)
    return slide


def slide_metrics(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.8), Inches(0.7))
    tf = tb.text_frame
    tf.text = "Результаты на отложенном тесте"
    style_title(tf, size=32)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.8), Inches(0.4))
    note.text_frame.text = "11 пар · модель не видела test.csv при обучении"
    note.text_frame.paragraphs[0].font.color.rgb = MUTED
    note.text_frame.paragraphs[0].font.size = Pt(13)

    cards = [
        ("85,71 %", "Precision\n(pred ≥ 0.5)"),
        ("0,971", "NDCG\nранжирование"),
        ("0,269", "Корреляция\nс экспертом"),
        ("206", "Слов\nв словаре"),
    ]
    xs = [0.7, 2.55, 4.4, 6.25]
    for x, (val, label) in zip(xs, cards):
        add_card(slide, x, 1.65, 1.65, 1.35, val, label, accent=ACCENT2 if val == "85,71 %" else ACCENT)

    img_path = GRAPHS / "precision_vs_threshold.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.7), Inches(3.2), width=Inches(4.2))
    img2 = GRAPHS / "predicted_vs_actual.png"
    if img2.exists():
        slide.shapes.add_picture(str(img2), Inches(5.1), Inches(3.2), width=Inches(4.2))

    add_footer(slide)
    return slide


def slide_architecture(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.8), Inches(0.7))
    tf = tb.text_frame
    tf.text = "Архитектура системы"
    style_title(tf, size=32)

    steps = [
        ("1", "Текст\nзапроса", ACCENT2),
        ("2", "Препро-\nцессор", ACCENT),
        ("3", "TextEmbedding\nModel", ACCENT),
        ("4", "SQLite\n+ поиск", ACCENT2),
        ("5", "REST API\n+ UI", SUCCESS),
    ]
    x0 = 0.55
    for i, (num, label, color) in enumerate(steps):
        x = x0 + i * 1.85
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(2.2), Inches(0.55), Inches(0.55)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        ntb = slide.shapes.add_textbox(Inches(x), Inches(2.28), Inches(0.55), Inches(0.4))
        ntb.text_frame.text = num
        ntb.text_frame.paragraphs[0].font.size = Pt(16)
        ntb.text_frame.paragraphs[0].font.bold = True
        ntb.text_frame.paragraphs[0].font.color.rgb = BG
        ntb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.15), Inches(3.0), Inches(1.55), Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = SURFACE
        box.line.color.rgb = RGBColor(0x2A, 0x35, 0x50)

        lbl = slide.shapes.add_textbox(Inches(x - 0.05), Inches(3.15), Inches(1.35), Inches(0.9))
        lbl.text_frame.text = label
        lbl.text_frame.paragraphs[0].font.size = Pt(11)
        lbl.text_frame.paragraphs[0].font.color.rgb = TEXT
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.35), Inches(3.35), Inches(0.45), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    stack = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(8.8), Inches(2.2))
    tf = stack.text_frame
    tf.word_wrap = True
    lines = [
        "Модель: Embedding(128) → BiLSTM(192) → Attention → FC → L2-norm → вектор 128D",
        "Обучение: Triplet Loss · train.csv (42 пары) · ~1,58 млн параметров",
        "Поиск: косинусная схожесть · фильтр game_type · top_k результатов",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸  " + line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)

    add_footer(slide)
    return slide


def slide_comparison(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.8), Inches(0.7))
    tf = tb.text_frame
    tf.text = "Сравнение с baseline-методами"
    style_title(tf, size=32)

    img = GRAPHS / "all_methods_comparison.png"
    if not img.exists():
        img = PROJECT / "methods_comparison.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.9), Inches(1.3), width=Inches(8.3))

    cap = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(8.8), Inches(0.5))
    cap.text_frame.text = (
        "Собственная модель: компактнее трансформеров, обучена на игровом сленге, "
        "Precision 85,71 % на hold-out test"
    )
    cap.text_frame.paragraphs[0].font.size = Pt(12)
    cap.text_frame.paragraphs[0].font.color.rgb = MUTED

    add_footer(slide)
    return slide


def slide_title_page(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # декоративные круги
    for size, x, y, alpha in [(4.5, 7.2, -1.5, SURFACE), (2.5, -0.8, 5.5, SURFACE)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        c.fill.solid()
        c.fill.fore_color.rgb = alpha
        c.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT2
    bar.line.fill.background()

    badge = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(4), Inches(0.4))
    badge.text_frame.text = "◆  Учебная практика · КТбо3-3"
    badge.text_frame.paragraphs[0].font.size = Pt(12)
    badge.text_frame.paragraphs[0].font.color.rgb = ACCENT2

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.8), Inches(2.2))
    tf = title_box.text_frame
    tf.text = TITLE
    style_title(tf, size=34)
    p2 = tf.add_paragraph()
    p2.text = SUBTITLE
    p2.font.size = Pt(20)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(14)

    info = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(8), Inches(1.8))
    itf = info.text_frame
    lines = [
        "Студент: Черкашина Екатерина Витальевна",
        "Направление: 09.03.01 Информатика и вычислительная техника",
        "Южный федеральный университет · ИКТИБ · 2025",
    ]
    for i, line in enumerate(lines):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)

    add_footer(slide, "Teammate Finder — semantic matchmaking for gamers")
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title_page(prs)

    slide_content(
        prs,
        "Проблема и актуальность",
        [
            "70 %+ игроков сталкиваются с несовместимыми тиммейтами в командных играх",
            "Существующие LFG-сервисы подбирают по рейтингу и тегам, а не по стилю игры",
            "Игровой сленг («агр сап», «каррю на лейте») ломает классические NLP-методы",
            "Нужен семантический поиск по смыслу описания, а не по ключевым словам",
        ],
        subtitle="Dota 2 · CS:GO · Valorant",
    )

    slide_content(
        prs,
        "Цель и задачи",
        [
            "Цель: разработать систему поиска совместимых тиммейтов по текстовому описанию стиля",
            "Проанализировать психологию командной игры и существующие платформы",
            "Реализовать препроцессор игрового сленга (80+ терминов)",
            "Обучить собственную модель TextEmbeddingModel на размеченном датасете",
            "Построить REST API, базу данных и веб-интерфейс",
            "Оценить качество на отложенной тестовой выборке (hold-out test)",
        ],
    )

    slide_architecture(prs)

    slide_content(
        prs,
        "Датасет и обучение",
        [
            "gaming_compatibility_dataset.csv — 53 экспертно размеченные пары (score 0–1)",
            "Разделение: train.csv (42) + test.csv (11) — честная оценка",
            "Positive (≥ 0.7) и negative (≤ 0.45) → триплеты для Triplet Loss",
            "30 эпох · Adam · validation loss → best_model.pth",
            "Примеры: «агр сап, давлю лайн» ↔ «агрессивный саппорт, хараслю линию»",
        ],
        subtitle="Triplet Loss · margin = 0.25",
    )

    slide_content(
        prs,
        "Технологический стек",
        [
            "Backend: FastAPI + Uvicorn",
            "ML: PyTorch — Embedding + BiLSTM + Attention",
            "БД: SQLite (users, embedding, game_type)",
            "NLP: собственный препроцессор text_processor.py",
            "Frontend: teammate_finder.html (тёмный UI, карточки совместимости)",
            "Оценка: scikit-learn — Precision, NDCG, корреляция Пирсона",
        ],
    )

    slide_metrics(prs)
    slide_comparison(prs)

    slide_content(
        prs,
        "Демонстрация",
        [
            "POST /similarity/search — поиск по описанию + game_type + top_k",
            "GET /model/info — статус модели, размер словаря, размерность эмбеддинга",
            "Swagger UI: http://localhost:8000/docs",
            "Веб-интерфейс: teammate_finder.html — карточки с % совместимости",
            "Запуск: python -m model.train → uvicorn api.main:app → seed_database.py",
        ],
        subtitle="REST API + веб-интерфейс",
    )

    slide_content(
        prs,
        "Выводы",
        [
            "Реализован прототип семантического матчмейкинга с собственной нейросетью",
            "Модель учитывает игровой сленг через препроцессор и предметное обучение",
            "На test: Precision 85,71 %, NDCG 0,971 — сопоставимо с Sentence-BERT",
            "Модель компактная (~1,58M параметров) vs тяжёлые трансформеры",
            "Перспективы: больше данных, pgvector, обратная связь, игровые API",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide, width=10, height=0.08, top=3.5)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.8), Inches(1.5))
    tf = tb.text_frame
    tf.text = "Спасибо за внимание!"
    style_title(tf, size=44, color=ACCENT2)
    p2 = tf.add_paragraph()
    p2.text = "Готова ответить на вопросы"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT
    p2.space_before = Pt(16)
    p2.alignment = PP_ALIGN.LEFT

    add_footer(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Презентация сохранена: {path}")
